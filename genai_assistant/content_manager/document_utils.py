import os
from docx import Document
from pptx import Presentation
from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from genai_assistant.settings import (
    AZURE_FORM_RECOGNIZER_ENDPOINT,
    AZURE_FORM_RECOGNIZER_KEY,
)

def analyze_document(file_path):
    """Analyze document using Azure Form Recognizer prebuilt-read model."""
    client = DocumentAnalysisClient(
        endpoint=AZURE_FORM_RECOGNIZER_ENDPOINT,
        credential=AzureKeyCredential(AZURE_FORM_RECOGNIZER_KEY),
    )
    with open(file_path, "rb") as fd:
        poller = client.begin_analyze_document("prebuilt-read", fd)
        result = poller.result()

    extracted_text = ""
    for page in result.pages:
        for line in page.lines:
            extracted_text += line.content + "\n"
    return extracted_text


def extract_text_from_file(file_path):
    """
    Extract text from a file path (used by file watcher).
    """
    ext = os.path.splitext(file_path)[1].lower()

    try:
        if ext == ".txt":
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read().strip()

        elif ext in [".pdf", ".jpg", ".jpeg", ".png", ".tiff", ".xlsx", ".csv"]:
            return analyze_document(file_path)

        elif ext == ".docx":
            doc = Document(file_path)
            return "\n".join(para.text for para in doc.paragraphs)

        elif ext == ".pptx":
            prs = Presentation(file_path)
            return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

        else:
            print(f"⚠️ Skipped unsupported file type: {file_path}")
            return ""

    except Exception as e:
        print(f"❌ Error extracting text from {file_path}: {e}")
        return ""


def extract_text_from_document(uploaded_items):
    """
    Process a list of uploaded items with .file attribute.
    Used during manual uploads.
    """
    full_text = ""

    for item in uploaded_items:
        if not hasattr(item, "file") or not item.file:
            continue

        file_path = item.file.path
        text = extract_text_from_file(file_path)
        full_text += text + "\n\n"

    return full_text.strip()
