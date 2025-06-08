from django.shortcuts import render

# Create your views here.
from rest_framework import generics
from .models import UploadedContent
from .serializers import UploadedContentSerializer

from .forms import UploadedContentForm
from django.shortcuts import redirect

from django.http import JsonResponse
import datetime

import json 
from django.views.decorators.csrf import csrf_exempt
import os

from django.conf import settings
from django.core.files.storage import default_storage
import uuid

from django.shortcuts import redirect, get_object_or_404
from django.contrib import messages
from .models import UploadedContent

from .models import ChatMessage

from .utils import extract_dates_and_references_from_file

from azure.ai.formrecognizer import DocumentAnalysisClient
from azure.core.credentials import AzureKeyCredential
from dateutil import parser

from docx import Document
from pptx import Presentation


from django.db.models import Count
import requests

from .document_utils import extract_text_from_file
from .langchain_utils import get_langchain_chain, embed_texts, split_text_to_chunks, generate_ai_summary
from .azure_search_utils import index_document, delete_document

import psutil

#API View

class UploadedContentListCreateView(generics.ListCreateAPIView):
    queryset = UploadedContent.objects.all().order_by('-uploaded_at')
    serializer_class = UploadedContentSerializer

#UI View
    
def home(request):
    return render(request, 'content_manager/home.html', {
        'weather_api_key': settings.WEATHER_API_KEY
    })


# Page to upload and manage content
'''def upload_page(request):
    return render(request, 'content_manager/upload.html')'''

#Form Page
def upload_content(request):
    if request.method == 'POST':
        form = UploadedContentForm(request.POST, request.FILES)
        if form.is_valid():
            uploaded_instance = form.save()

            if uploaded_instance.file:
                file_path = uploaded_instance.file.path
                print(f"Extracting text from file at: {file_path}")
                extracted_text = extract_text_from_file(file_path)
                print(f"Extracted text length: {len(extracted_text)}")
                uploaded_instance.extracted_text = extracted_text
                uploaded_instance.save()

                # Split, embed, and index chunks
                chunks = split_text_to_chunks(extracted_text)
                print(f"Split into {len(chunks)} chunks")
                vectors = embed_texts(chunks)
                print(f"Generated {len(vectors)} vectors")

                vector_ids = []
                for chunk, vector in zip(chunks, vectors):
                    doc_id = index_document(uploaded_instance.title, chunk, vector)
                    vector_ids.append(doc_id)

                uploaded_instance.vector_id = ",".join(vector_ids)
                uploaded_instance.save()

            return redirect('content_list')
    else:
        form = UploadedContentForm()
    return render(request, 'content_manager/upload.html', {'form': form})

# Page to show all uploaded content
def content_list_page(request):
    uploaded_contents = UploadedContent.objects.all().order_by('-uploaded_at')
    return render(request, 'content_manager/content_list.html', {'uploaded_contents': uploaded_contents})

# Chatbot interface page
def chatbot_page(request):
    return render(request, 'content_manager/chatbot.html')

# Calendar 
'''def calendar_events(request):
    events = [
        {
            "title": "Team Review",
            "start": str(datetime.date.today()),
        },
        {
            "title": "Deadline - Project Plan",
            "start": str(datetime.date.today() + datetime.timedelta(days=3)),
        },
    ]
    return JsonResponse(events, safe=False)'''

import re
def remove_dates(text):
    date_patterns = [
        r'\b\d{1,2}[-/]\d{1,2}[-/]\d{2,4}\b',      # 12/05/2024 or 12-05-2024
        r'\b\d{4}[-/]\d{1,2}[-/]\d{1,2}\b',        # 2024/05/12
        r'\b(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{1,2},?\s+\d{4}\b', # Jan 1, 2024
        r'\b\d{1,2}\s+(?:Jan|Feb|Mar|Apr|May|Jun|Jul|Aug|Sep|Oct|Nov|Dec)[a-z]*\s+\d{4}\b'    # 1 Jan 2024
    ]
    for pattern in date_patterns:
        text = re.sub(pattern, '', text, flags=re.IGNORECASE)
    return text.strip()

def extract_dates_from_text(text):
    from dateutil.parser import parse
    events = []
    lines = text.split('\n')

    for line in lines:
        try:
            date = parse(line, fuzzy=True)

            # Determine category based on keywords
            line_lower = line.lower()
            if any(keyword in line_lower for keyword in ["meeting", "review", "project", "deadline", "strategy","risk","test","develop"]):
                category = "professional"
                color = "#87CEFA"  # Light Blue
            elif any(keyword in line_lower for keyword in ["birthday", "anniversary", "party", "doctor", "appointment", "family", "friend"]):
                category = "personal"
                color = "#FFB6C1"  # Light Pink
            else:
                category = "general"
                color = "#D3D3D3"  # Light Gray for uncategorized

            events.append({
                "title": remove_dates(line)[:50],
                "start": date.date().isoformat(),
                "category": category,
                "color": color
            })

        except:
            continue

    return events


def analyze_document(file_path):
    endpoint = settings.AZURE_DOCUMENT_INTELLIGENCE_ENDPOINT
    key = settings.AZURE_DOCUMENT_INTELLIGENCE_KEY
    client = DocumentAnalysisClient(endpoint=endpoint, credential=AzureKeyCredential(key))

    with open(file_path, "rb") as f:
        poller = client.begin_analyze_document("prebuilt-read", document=f)
        result = poller.result()

    content = ""
    for page in result.pages:
        for line in page.lines:
            content += line.content + "\n"
    return content

def calendar_events(request):
    events = []
    uploaded = UploadedContent.objects.all()

    for item in uploaded:
        if item.file:
            file_path = item.file.path
            ext = os.path.splitext(file_path)[1].lower()

            try:
                if ext == ".txt":
                    # Direct read for .txt
                    with open(file_path, "r", encoding="utf-8") as f:
                        text = f.read()
                elif ext in [".pdf", ".jpg", ".jpeg", ".png", ".tiff", ".xlsx"]:
                    text = analyze_document(file_path)
                elif ext == ".docx":
                    doc = Document(file_path)
                    text = "\n".join([para.text for para in doc.paragraphs])
                elif ext == ".pptx":
                    prs = Presentation(file_path)
                    text_runs = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                text_runs.append(shape.text)
                    text = "\n".join(text_runs)
                else:
                    print(f"⚠️ Skipped unsupported file type: {file_path}")
                    continue

                extracted = extract_dates_from_text(text)
                events.extend(extracted)

            except Exception as e:
                print(f"❌ Error processing {item.title}: {e}")
                continue

    return JsonResponse(events, safe=False)

    

def calendar_test_view(request):
    return render(request, 'content_manager/calendar_test.html')


# Chatbot API
from .langchain_utils import get_langchain_chain  
from .urls_utils import extract_text_from_url, save_text_as_file

@csrf_exempt
def chatbot_api(request):
    if request.method != "POST":
        return JsonResponse({"error": "Only POST requests are allowed."}, status=405)

    message = request.POST.get("message", "").strip()
    uploaded_file = request.FILES.get("file")
    ai_reply = ""

    # Check if message contains a URL
    url_pattern = r'(https?://[^\s]+)'
    urls = re.findall(url_pattern, message)

    if urls:
        url = urls[0]
        print(f"Detected URL: {url}")
        extracted_text = extract_text_from_url(url)
        print(f"Extracted text length: {len(extracted_text)}")

        if extracted_text.startswith("Error extracting text:"):
            ai_reply += extracted_text + "\n"
        else:
            filename = save_text_as_file(extracted_text, prefix="scraped_")
            print(f"Saved file: {filename}")

            # Save DB entry
            uploaded_instance = UploadedContent.objects.create(
                title=f"Scraped content from {url}",
                file=filename,
                content_type='url',
                extracted_text=extracted_text,
            )

            # Chunk, embed, index
            try:
                chunks = split_text_to_chunks(extracted_text)
                print(f"Chunks generated: {len(chunks)}")

                vectors = embed_texts(chunks)
                print(f"Vectors generated: {len(vectors)}")
            except Exception as e:
                print(f"Error during chunking/embedding: {e}")
                ai_reply += f"❌ Error during processing: {str(e)}"
                chunks = []
                vectors = []

            vector_ids = []
            for chunk, vector in zip(chunks, vectors):
                try:
                    doc_id = index_document(filename, chunk, vector)
                    vector_ids.append(doc_id)
                    print(f"Indexed chunk doc_id: {doc_id}")
                except Exception as e:
                    print(f"Error indexing document chunk: {e}")

            uploaded_instance.vector_id = ",".join(vector_ids)
            uploaded_instance.save()

            ai_reply += f"✅ Fetched and processed content from URL: {url}\n"

    # Handle uploaded file (normal upload)
    elif uploaded_file:
        filename = f"{uuid.uuid4()}_{uploaded_file.name}"
        file_path = os.path.join("uploads", filename)
        save_path = os.path.join(settings.MEDIA_ROOT, file_path)
        os.makedirs(os.path.dirname(save_path), exist_ok=True)

        with default_storage.open(save_path, 'wb+') as destination:
            for chunk in uploaded_file.chunks():
                destination.write(chunk)

        print(f"Extracting text from uploaded file at: {save_path}")
        extracted_text = extract_text_from_file(save_path)

        uploaded_instance = UploadedContent.objects.create(
            title=uploaded_file.name,
            file=file_path,
            content_type='doc',
            extracted_text=extracted_text,
        )

        try:
            chunks = split_text_to_chunks(extracted_text)
            print(f"Chunks generated: {len(chunks)}")

            vectors = embed_texts(chunks)
            print(f"Vectors generated: {len(vectors)}")
        except Exception as e:
            print(f"Error during chunking/embedding: {e}")
            ai_reply += f"❌ Error during processing: {str(e)}"
            chunks = []
            vectors = []

        vector_ids = []

        for chunk, vector in zip(chunks, vectors):
            try:
                doc_id = index_document(uploaded_file.name, chunk, vector)
                vector_ids.append(doc_id)
                print(f"Indexed chunk doc_id: {doc_id}")
            except Exception as e:
                print(f"Error indexing document chunk: {e}")

        uploaded_instance.vector_id = ",".join(vector_ids)
        uploaded_instance.save()

        ai_reply += f"✅ Uploaded and processed file: {uploaded_file.name}. "

    # Handle chat message (non-URL, non-file)
    if message and not urls:
        try:
            rag_function = get_langchain_chain()  # returns a callable
            langchain_response = rag_function(message)
            ai_reply += langchain_response
        except Exception as e:
            ai_reply += f"❌ Error generating response: {str(e)}"
    elif not ai_reply:
        ai_reply = "🤖 You didn't say anything!"

    # Save chat history
    ChatMessage.objects.create(
        user_message=message or "",
        bot_reply=ai_reply
    )

    return JsonResponse({"reply": ai_reply})




    
#Delete Saved Files
def delete_content(request, content_id):
    content = get_object_or_404(UploadedContent, id=content_id)

    if request.method == "POST":
        file_path = content.file.path if content.file else None

        # Delete file safely
        if file_path and os.path.isfile(file_path):
            try:
                os.remove(file_path)
            except PermissionError as e:
                # Check which process is locking the file (Windows)
                for proc in psutil.process_iter(["pid", "name", "open_files"]):
                    try:
                        for f in proc.info["open_files"] or []:
                            if file_path == f.path:
                                print(f"⚠️ File in use by PID {proc.info['pid']} - {proc.info['name']}")
                    except Exception:
                        pass
                messages.error(request, f"File is currently in use and cannot be deleted: {e}")
                return redirect('content_list')
            except Exception as e:
                messages.error(request, f"Unexpected error deleting file: {e}")
                return redirect('content_list')

        # Delete from Azure AI Search index
        if content.vector_id:
            try:
                for doc_id in content.vector_id.split(","):
                    delete_document(doc_id.strip())
            except Exception as e:
                messages.warning(request, "File deleted, but some vector data may not have been removed.")

        content.delete()
        messages.success(request, "File and associated data deleted successfully.")

    return redirect('content_list')

'''
#Save chats
def chatbot_api(request):
    if request.method == 'POST':
        data = json.loads(request.body)
        user_message = data.get('message')
        # Simulate response (replace with real logic)
        bot_reply = f"Echo: {user_message}"

        # Save chat
        ChatMessage.objects.create(user_message=user_message, bot_reply=bot_reply)

        return JsonResponse({'reply': bot_reply})'''

   
#Save Chats Page
def saved_chats(request):
    chats = ChatMessage.objects.order_by('-timestamp')
    return render(request, 'content_manager/saved_chats.html', {'chats': chats})




# Explore Section
from .langchain_utils import get_langchain_chain_updated  

@csrf_exempt
def explore_api(request):
    if request.method == "POST":
        query = request.POST.get("query", "").strip()

        if not query:
            return JsonResponse({"error": "No input provided."}, status=400)

        try:
            
            chain = get_langchain_chain_updated()
            response = chain.run(query)

            return JsonResponse({"response": response})

        except Exception as e:
            return JsonResponse({"error": str(e)}, status=500)

    return JsonResponse({"error": "Only POST allowed."}, status=405)
#explore.html
def explore_page(request):
    return render(request, 'content_manager/explore.html')

#Promt generation - RAG based
def ai_insights(request):
    """
    Generate AI-based insights from Azure AI Search indexed data.
    Includes: Pending Work Items and Upcoming Workloads.
    """
    try:
        # Use LLM summaries directly based on embedded vector data
        pending_prompt = "List 3 pending or incomplete work items from the documents with priority category. Each item of the list should have title, description, date and priority."
        upcoming_prompt = "Predict 3 upcoming workloads or expected tasks based on the documents."

        pending_summary = generate_ai_summary(pending_prompt)
        upcoming_summary = generate_ai_summary(upcoming_prompt)

        return JsonResponse({
            "pending": [pending_summary],
            "upcoming": [upcoming_summary]
        })

    except Exception as e:
        return JsonResponse({"error": str(e)}, status=500)